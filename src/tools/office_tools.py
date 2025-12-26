"""Office document operations tools"""

import logging
from pathlib import Path
from typing import Dict, Any

from .security import validate_path

logger = logging.getLogger('office_tools')

# Try to import optional dependencies
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False


def _error_response(error: str, error_type: str) -> Dict[str, Any]:
    """Create standardized error response"""
    return {
        'success': False,
        'result': None,
        'error': error,
        'error_type': error_type
    }


def _success_response(result: Any) -> Dict[str, Any]:
    """Create standardized success response"""
    return {
        'success': True,
        'result': result,
        'error': None,
        'error_type': None
    }


def docx_extract_text(docx_path: str) -> Dict[str, Any]:
    """
    Extract text from DOCX
    
    Args:
        docx_path: Path to DOCX file
        
    Returns:
        Dict with success, result (extracted text), error, error_type
    """
    if not DOCX_AVAILABLE:
        return _error_response("python-docx not available. Install with: pip install python-docx", 'ImportError')
    
    try:
        is_valid, error_msg = validate_path(docx_path, operation='read')
        if not is_valid:
            return _error_response(error_msg, 'ValidationError')
        
        path_obj = Path(docx_path)
        doc = Document(path_obj)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            text_parts.append(paragraph.text)
        
        text = '\n'.join(text_parts)
        return _success_response(text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {docx_path}: {e}", exc_info=True)
        return _error_response(str(e), type(e).__name__)


def docx_get_metadata(docx_path: str) -> Dict[str, Any]:
    """
    Get DOCX metadata
    
    Args:
        docx_path: Path to DOCX file
        
    Returns:
        Dict with success, result (metadata dict), error, error_type
    """
    if not DOCX_AVAILABLE:
        return _error_response("python-docx not available. Install with: pip install python-docx", 'ImportError')
    
    try:
        is_valid, error_msg = validate_path(docx_path, operation='read')
        if not is_valid:
            return _error_response(error_msg, 'ValidationError')
        
        path_obj = Path(docx_path)
        doc = Document(path_obj)
        
        core_props = doc.core_properties
        
        metadata = {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
            'keywords': core_props.keywords or '',
            'created': str(core_props.created) if core_props.created else '',
            'modified': str(core_props.modified) if core_props.modified else '',
        }
        
        return _success_response(metadata)
    except Exception as e:
        logger.error(f"Error getting DOCX metadata {docx_path}: {e}", exc_info=True)
        return _error_response(str(e), type(e).__name__)


def pptx_extract_text(pptx_path: str) -> Dict[str, Any]:
    """
    Extract text from PPTX
    
    Args:
        pptx_path: Path to PPTX file
        
    Returns:
        Dict with success, result (extracted text), error, error_type
    """
    if not PPTX_AVAILABLE:
        return _error_response("python-pptx not available. Install with: pip install python-pptx", 'ImportError')
    
    try:
        is_valid, error_msg = validate_path(pptx_path, operation='read')
        if not is_valid:
            return _error_response(error_msg, 'ValidationError')
        
        path_obj = Path(pptx_path)
        prs = Presentation(path_obj)
        
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        text = '\n\n'.join(text_parts)
        return _success_response(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {pptx_path}: {e}", exc_info=True)
        return _error_response(str(e), type(e).__name__)


def xlsx_read_all_sheets(xlsx_path: str) -> Dict[str, Any]:
    """
    Read all Excel sheets
    
    Args:
        xlsx_path: Path to XLSX file
        
    Returns:
        Dict with success, result (dict of sheet data), error, error_type
    """
    if not XLSX_AVAILABLE:
        return _error_response("openpyxl not available. Install with: pip install openpyxl", 'ImportError')
    
    try:
        is_valid, error_msg = validate_path(xlsx_path, operation='read')
        if not is_valid:
            return _error_response(error_msg, 'ValidationError')
        
        path_obj = Path(xlsx_path)
        wb = openpyxl.load_workbook(path_obj, data_only=True)
        
        result = {}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append(list(row))
            result[sheet_name] = rows
        
        return _success_response(result)
    except Exception as e:
        logger.error(f"Error reading XLSX {xlsx_path}: {e}", exc_info=True)
        return _error_response(str(e), type(e).__name__)


def xlsx_get_sheet_names(xlsx_path: str) -> Dict[str, Any]:
    """
    Get sheet names
    
    Args:
        xlsx_path: Path to XLSX file
        
    Returns:
        Dict with success, result (list of sheet names), error, error_type
    """
    if not XLSX_AVAILABLE:
        return _error_response("openpyxl not available. Install with: pip install openpyxl", 'ImportError')
    
    try:
        is_valid, error_msg = validate_path(xlsx_path, operation='read')
        if not is_valid:
            return _error_response(error_msg, 'ValidationError')
        
        path_obj = Path(xlsx_path)
        wb = openpyxl.load_workbook(path_obj)
        
        return _success_response(wb.sheetnames)
    except Exception as e:
        logger.error(f"Error getting sheet names {xlsx_path}: {e}", exc_info=True)
        return _error_response(str(e), type(e).__name__)

